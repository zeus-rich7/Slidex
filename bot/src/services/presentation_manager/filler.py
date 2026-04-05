from itertools import zip_longest

from pptx import Presentation

from src.helpers.async_factory import to_async


class Slide:
    def __init__(self, title: str = None, author: str = None, recipient: str = None,
                 lang: str = None, slides: int = None, chars_per_content: int = None,
                 template_path: str = None, output_path: str = None, json_data = None):
        self.title = title
        self.author = author
        self.recipient = recipient
        self.lang = lang
        self.slides = slides
        self.chars_per_content: int = chars_per_content
        self.template_path = template_path
        self.json_data = json_data
        self.__presentation = Presentation(template_path)
        self.__output_path: str = output_path


    def add_title_slide(self) -> bool:
        try:
            title_layout = self.__presentation.slide_layouts[0]
            self.__presentation.slides.add_slide(slide_layout=title_layout)

            title_page = self.__presentation.slides[-1]

            for x,y in zip(title_layout.placeholders, title_page.placeholders):
                y.name = x.name

            for placeholder in title_page.placeholders:
                if placeholder.name == "title":
                    placeholder.text = self.title
                elif placeholder.name == "author":
                    placeholder.text = self.author
                elif placeholder.name == "recipient":
                    placeholder.text = self.recipient
            return True
        except IndexError:
            return False

    def add_plans_slide(self)-> bool:
        try:
            plans_layout = self.__presentation.slide_layouts[1]
            self.__presentation.slides.add_slide(slide_layout=plans_layout)

            plans_page = self.__presentation.slides[-1]

            for x,y in zip(plans_layout.placeholders, plans_page.placeholders):
                y.name = x.name

            for placeholder, content in zip_longest(plans_page.placeholders, self.json_data[0]["contents"]):
                if content:
                    if placeholder.name == "plan":
                        placeholder.text = self.json_data[0]["header"]
                    else:
                        placeholder.text = content["title"]
                else:
                    placeholder.element.getparent().remove(placeholder.element)
            return True
        except IndexError:
            return False

    def add_content_slides(self):
        for slide_layout, json_layout in zip(list(self.__presentation.slide_layouts)[2:], self.json_data[1:]):
            self.__presentation.slides.add_slide(slide_layout)

            page = self.__presentation.slides[-1]

            for x,y in zip(page.placeholders, slide_layout.placeholders):
                x.name = y.name

            for placeholder, content in zip(list(page.placeholders), json_layout["contents"]):
                if placeholder.name == "title":
                    placeholder.text = json_layout["header"]
                if placeholder.name.startswith("subtitle"):
                    placeholder.text = content["title"]
                elif placeholder.name.startswith("content"):
                    placeholder.text = content["content"]


    @to_async()
    def save(self) -> str:
        self.__presentation.save(self.__output_path)
        return self.__output_path

    @property
    def output_path(self):
        return self.__output_path

    @output_path.setter
    def output_path(self, val):
        self.__output_path = val

    @to_async()
    def construct(self):
        self.add_title_slide()
        self.add_plans_slide()
        self.add_content_slides()














